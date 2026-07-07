#!/usr/bin/env python3
"""make_ppt.py v4  —  웹 테마(라이트 블루/화이트) PPT, 풀 matplotlib 렌더링"""
import io, os
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.font_manager as fm
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from matplotlib.gridspec import GridSpec
import matplotlib.patheffects as pe
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ─────────────────────────────────────────────────────────────────────────────
#  Korean font
# ─────────────────────────────────────────────────────────────────────────────
import sys as _sys
if _sys.platform == 'win32':
    _KR  = r'C:\Windows\Fonts\malgun.ttf'
    _KRB = r'C:\Windows\Fonts\malgunbd.ttf'
else:
    _KR  = '/mnt/c/Windows/Fonts/malgun.ttf'
    _KRB = '/mnt/c/Windows/Fonts/malgunbd.ttf'
if os.path.exists(_KR):
    fm.fontManager.addfont(_KR)
if os.path.exists(_KRB):
    fm.fontManager.addfont(_KRB)
matplotlib.rcParams['font.family'] = 'Malgun Gothic'
matplotlib.rcParams['axes.unicode_minus'] = False

def _font(bold=False):
    return {'fontproperties': fm.FontProperties(fname=_KRB if bold else _KR)}

# ─────────────────────────────────────────────────────────────────────────────
#  Color palette  (웹 CSS 변수 그대로)
# ─────────────────────────────────────────────────────────────────────────────
P = {
    'bg':   '#FFFFFF',
    'bg2':  '#F5F7FA',
    'bg3':  '#EBF4FC',
    'bg4':  '#E0F4FF',
    'bg5':  '#B3E0F7',
    'pri':  '#0076CE',
    'pdk':  '#005BA3',
    'sky':  '#5BB8F5',
    'grn':  '#2FAA6E',
    'prp':  '#5B7FD6',
    'amb':  '#F0A830',
    'red':  '#E05050',
    'txt':  '#1A1A2E',
    'dim':  '#4A6080',
    'mute': '#8BA5BE',
    'brd':  '#E0EAF4',
    'brh':  '#B8D0E8',
}

def hx(h):
    if isinstance(h, tuple): return h
    if h == 'white': return (1,1,1)
    if h == 'black': return (0,0,0)
    return tuple(int(h[i:i+2], 16)/255 for i in (1,3,5))

# ─────────────────────────────────────────────────────────────────────────────
#  Shared rendering helpers
# ─────────────────────────────────────────────────────────────────────────────

def gradient_bg(ax, W=1, H=1, alpha=1.0):
    """라이트 그라데이션 배경  (왼쪽 위=흰색, 오른쪽 아래=하늘색)"""
    xi = np.linspace(0, 1, 200)[None, :]
    yi = np.linspace(0, 1, 200)[:, None]
    t  = (0.5*xi + 0.5*yi)
    r  = 1.0 - t*(1.0 - 179/255)
    g  = 1.0 - t*(1.0 - 224/255)
    b  = 1.0 - t*(1.0 - 247/255)
    img = np.stack([r,g,b], axis=-1)
    ax.imshow(img, extent=[0,W,0,H], aspect='auto', origin='upper',
              zorder=0, alpha=alpha)

def grid_overlay(ax, W, H, spacing=0.6, alpha=0.04):
    col = (*hx(P['pri']), alpha)
    for x in np.arange(0, W+spacing, spacing):
        ax.plot([x,x],[0,H], color=col, lw=0.5, zorder=1)
    for y in np.arange(0, H+spacing, spacing):
        ax.plot([0,W],[y,y], color=col, lw=0.5, zorder=1)

def pill(ax, x, y, w, h, color, radius=0.1, alpha=1.0, zorder=3):
    p = FancyBboxPatch((x,y), w, h,
                       boxstyle=f'round,pad={radius}',
                       facecolor=hx(color), edgecolor='none',
                       alpha=alpha, zorder=zorder)
    ax.add_patch(p)
    return p

def card(ax, x, y, w, h, fc=None, ec=None, lw=1.2, shadow=True, zorder=3):
    if shadow:
        ax.add_patch(FancyBboxPatch((x+0.04, y-0.04), w, h,
                                    boxstyle='round,pad=0.08',
                                    facecolor=hx(P['brd']), edgecolor='none',
                                    alpha=0.8, zorder=zorder-1))
    ax.add_patch(FancyBboxPatch((x,y), w, h,
                                boxstyle='round,pad=0.08',
                                facecolor=fc or hx(P['bg']),
                                edgecolor=ec or hx(P['brh']),
                                linewidth=lw, zorder=zorder))

def arrow(ax, x1, y1, x2, y2, color=None, lw=1.8, ms=14):
    ax.annotate('', xy=(x2,y2), xytext=(x1,y1),
                arrowprops=dict(arrowstyle='->', color=color or hx(P['pri']),
                                lw=lw, mutation_scale=ms))

def accent_bar(ax, x, y, h, color, w=0.06, zorder=5):
    ax.add_patch(mpatches.Rectangle((x,y), w, h,
                                    facecolor=hx(color), edgecolor='none', zorder=zorder))

def fig2stream(fig, dpi=144, bg='white'):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight',
                facecolor=bg, edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf

# text with Korean font shorthand
def T(ax, x, y, s, size=11, color=None, bold=False, ha='left', va='center',
      alpha=1.0, zorder=10, **kw):
    fp = fm.FontProperties(fname=_KRB if bold else _KR)
    c  = color or P['txt']
    c  = hx(c) if isinstance(c, str) and c.startswith('#') else c
    ax.text(x, y, s, fontsize=size, color=c,
            ha=ha, va=va, fontproperties=fp, alpha=alpha, zorder=zorder, **kw)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 1  –  타이틀
# ─────────────────────────────────────────────────────────────────────────────
def slide_title():
    fig = plt.figure(figsize=(13.33, 7.5), facecolor='white')
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0, 13.33); ax.set_ylim(0, 7.5); ax.axis('off')

    gradient_bg(ax, W=13.33, H=7.5)
    grid_overlay(ax, W=13.33, H=7.5, spacing=0.72, alpha=0.035)

    # 왼쪽 파란 세로선
    ax.add_patch(mpatches.Rectangle((0,0), 0.08, 7.5,
                                    facecolor=hx(P['pri']), edgecolor='none', zorder=4))

    # 오른쪽 패널 — 악보 블록 데모
    card(ax, 7.1, 0.5, 5.9, 6.5, fc=hx(P['bg']), shadow=True, zorder=2)
    T(ax, 10.05, 6.55, '커스텀 표기법 미리보기', size=10, color=P['mute'], ha='center')

    notes   = ['C','D','E','F','G','A','B']
    n_col   = [P['pri'],P['sky'],P['grn'],'#F5A623',P['amb'],P['prp'],P['red']]
    dur     = [1.0, 0.5, 0.5, 1.0, 0.5, 0.5, 1.0]
    kor     = ['도','레','미','파','솔','라','시']
    xn = 7.35
    for n,nc,d,k in zip(notes, n_col, dur, kor):
        bw = d*0.79
        pill(ax, xn, 5.42, bw-0.06, 0.95, nc, radius=0.07)
        T(ax, xn+bw/2-0.03, 5.92, n, size=17, color=P['bg'], bold=True, ha='center')
        T(ax, xn+bw/2-0.03, 5.52, k, size=9,  color=P['bg'], ha='center')
        xn += bw

    # 오선 장식선
    for yy in [4.95, 5.09, 5.23, 5.37]:
        ax.plot([7.3, 12.95], [yy, yy], color=hx(P['brh']), lw=0.9, zorder=3)

    # 박자 블록 예시
    T(ax, 10.05, 4.72, '박자 비례 너비', size=9, color=P['mute'], ha='center')
    dur2 = [('C',1.0,P['pri']),('E',0.5,P['grn']),('E',0.5,P['grn']),
            ('G',1.0,P['amb']),('A',2.0,P['sky'])]
    xd = 7.35; tw = 5.55; unit = tw / 5
    for nn, d2, nc2 in dur2:
        bw2 = d2*unit - 0.05
        pill(ax, xd, 4.0, bw2, 0.62, nc2, radius=0.06)
        T(ax, xd+bw2/2, 4.31, nn, size=14, color=P['bg'], bold=True, ha='center')
        xd += d2*unit

    # 피아노 건반 (간략)
    T(ax, 10.05, 3.67, '건반 배치', size=9, color=P['mute'], ha='center')
    n_w=14; kw_=5.4/n_w; x0=7.35
    bk={1,2,4,5,6,8,9,11,12}
    nc3=[P['pri'],P['sky'],P['grn'],'#F5A623',P['amb'],P['prp'],P['red']]*2
    for i in range(n_w):
        pill(ax, x0+i*kw_+0.01, 2.62, kw_-0.04, 0.95,
             P['bg2'] if i not in bk else P['bg'], radius=0.04, zorder=4)
        ax.add_patch(mpatches.Rectangle((x0+i*kw_+0.01, 2.62), kw_-0.04, 0.06,
                                        facecolor=hx(nc3[i]), edgecolor='none', zorder=5))
    for i in bk:
        if i < n_w:
            pill(ax, x0+i*kw_+kw_*0.55, 3.0, kw_*0.45, 0.57,
                 P['txt'], radius=0.03, zorder=6)

    T(ax, 10.05, 2.38, '흰 건반 = 영문자  /  검은 건반 = 숫자', size=8.5,
      color=P['mute'], ha='center')

    # 왼쪽 — 타이틀 텍스트
    T(ax, 0.5,  6.68, 'AI OMR 악보 변환 앱', size=10.5, color=P['pri'], bold=True)
    T(ax, 0.5,  5.8,  '맞춤형 악보 인식', size=42, bold=True, color=P['txt'], va='center')
    T(ax, 0.5,  4.95, '& 변환 앱', size=42, bold=True, color=P['pri'], va='center')
    ax.plot([0.5, 3.2], [4.42, 4.42], color=hx(P['pri']), lw=2.5)
    T(ax, 0.5, 4.05, '누구나 음악에 "연결"될 수 있도록', size=15, color=P['dim'])
    T(ax, 0.5, 3.55, '악보 사진 → AI 인식 → 쉬운 표기법 → 연주 연습', size=11.5, color=P['mute'])

    # 진행 현황 뱃지들
    badges = [
        ('Round 1  84–92%', P['grn']),
        ('Round 2  재학습 중', P['amb']),
        ('Flutter 앱 UI  완료', P['pri']),
        ('데모  2026·08·17', P['dim']),
    ]
    yb = 2.65
    for txt, col in badges:
        pill(ax, 0.5, yb, 5.8, 0.46, col, radius=0.08, alpha=0.1, zorder=3)
        ax.add_patch(mpatches.Rectangle((0.5, yb), 0.05, 0.46,
                                        facecolor=hx(col), edgecolor='none', zorder=4))
        T(ax, 0.7, yb+0.23, txt, size=11.5, color=col, bold=True)
        yb -= 0.57

    T(ax, 0.5, 0.35, '한양대학교 / 2026 캡스톤 디자인', size=9, color=P['mute'])
    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2  –  프로젝트 개요
# ─────────────────────────────────────────────────────────────────────────────
def slide_overview():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.35)

    # 헤더
    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '프로젝트 소개', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '악보를 못 읽어도 OK — 직관적 기호 변환으로 누구나 음악을', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '03 / 16', size=9, color=P['mute'], ha='right')

    # ── 왼쪽 패널: 일반 악보 ─────────────────────────────────────────────
    lx, lw = 0.35, 5.6
    card(ax, lx, 0.35, lw, 5.75, fc=(1,1,1), shadow=True, lw=1.2)
    ax.add_patch(mpatches.Rectangle((lx, 5.72), lw, 0.38,
                 facecolor=hx(P['mute']), edgecolor='none', zorder=5))
    T(ax, lx+lw/2, 5.91, '일반 악보  (Standard Notation)', size=11, bold=True,
      color='white', ha='center', zorder=6)

    # 오선지 두 단 그리기
    for sy, clef in [(4.52, '𝄞'), (3.05, 'F')]:
        for li in range(5):
            yline = sy - 0.32 + li*0.16
            ax.plot([lx+0.22, lx+lw-0.2], [yline, yline],
                    color='#1a1a1a', lw=0.85, zorder=5)
        # 세로 줄
        for xbar in [lx+0.22, lx+2.25, lx+lw-0.2]:
            ax.plot([xbar,xbar],[sy-0.32, sy+0.32], color='#1a1a1a', lw=1.1, zorder=5)
        # 음자리표 레이블
        clef_lbl = 'G' if clef == '𝄞' else 'F'
        T(ax, lx+0.46, sy+0.06, clef_lbl, size=18, bold=True,
          color='#333333', ha='center', zorder=6)
        # 음표
        for xo, yo in [(0.85,-0.16),(1.12,0.0),(1.38,0.16),(1.65,-0.08),
                       (1.9,0.08),(2.5,-0.16),(2.78,0.0),(3.05,0.16)]:
            nx, ny = lx+xo, sy+yo
            el = mpatches.Ellipse((nx, ny), 0.14, 0.10, angle=-22,
                                  facecolor='#111', edgecolor='none', zorder=6)
            ax.add_patch(el)
            ax.plot([nx+0.07, nx+0.07],[ny, ny+0.38],
                    color='#111', lw=0.8, zorder=6)

    # 악보 어려움 태그
    ax.plot([lx+0.22, lx+lw-0.2],[2.35, 2.35], color=hx(P['brd']), lw=0.6)
    for ki, txt in enumerate(['오선·음표 해독에 전문 교육 필요',
                               '기호 위치·길이를 모두 암기해야 함',
                               '악보 파일 없이는 앱 사용 불가']):
        py = 2.1 - ki*0.5
        ax.plot([lx+0.36],[py], 's', ms=4.5, color=hx(P['red']), zorder=6)
        T(ax, lx+0.5, py, txt, size=10, color=P['red'], zorder=6)

    # ── 중앙 OMR 변환 ─────────────────────────────────────────────────────
    mx = 6.67
    arrow(ax, lx+lw+0.06, 3.8, mx-0.62, 3.8, color=P['pri'], lw=2.2, ms=16)
    pill(ax, mx-0.62, 3.42, 1.24, 0.72, P['pri'], radius=0.1, zorder=6)
    T(ax, mx, 3.84, 'OMR', size=13, bold=True, color='white', ha='center', zorder=7)
    T(ax, mx, 3.54, 'AI 인식', size=9.5, color='white', ha='center', zorder=7)
    arrow(ax, mx+0.62, 3.8, 7.38+0.08, 3.8, color=P['grn'], lw=2.2, ms=16)

    # ── 오른쪽 패널: 커스텀 표기법 ────────────────────────────────────────
    rx, rw = 7.38, 5.6
    card(ax, rx, 0.35, rw, 5.75, fc=(1,1,1), shadow=True, lw=2.0)
    ax.add_patch(mpatches.Rectangle((rx, 5.72), rw, 0.38,
                 facecolor=hx(P['pri']), edgecolor='none', zorder=5))
    T(ax, rx+rw/2, 5.91, '커스텀 표기법  (AndroMR)', size=11, bold=True,
      color='white', ha='center', zorder=6)

    # 커스텀 기호 블록 두 행
    rows = [
        [('도','C',P['pri']),('레','D',P['sky']),('미','E',P['grn']),('파','F',P['amb'])],
        [('솔','G',P['prp']),('라','A',P['red']),('시','B',P['pri']),('도','C',P['sky'])],
    ]
    for ri, (row, sy) in enumerate(zip(rows, [4.55, 3.35])):
        for ci, (kor, eng, col) in enumerate(row):
            sx = rx + 0.62 + ci*1.2
            circ = plt.Circle((sx, sy), 0.33, facecolor=hx(col),
                               edgecolor='white', linewidth=2, zorder=6)
            ax.add_patch(circ)
            T(ax, sx, sy+0.08, kor, size=11.5, bold=True, color='white',
              ha='center', zorder=7)
            T(ax, sx, sy-0.14, eng, size=8, color='white', ha='center', zorder=7)
            if ci < 3:
                ax.plot([sx+0.33, sx+0.87],[sy, sy], color=hx(P['brd']),
                        lw=1.0, zorder=5)

    # 공유 기능 행
    ax.plot([rx+0.22, rx+rw-0.2],[2.68, 2.68], color=hx(P['brd']), lw=0.6)
    T(ax, rx+rw/2, 2.48, '앱에서 바로 연주 & 공유', size=11, bold=True,
      color=P['pri'], ha='center', zorder=5)
    for si, (lbl, col) in enumerate([('연주', P['grn']),('녹음', P['sky']),('공유', P['prp'])]):
        sx2 = rx+0.85+si*1.72
        pill(ax, sx2, 1.9, 1.38, 0.42, col, radius=0.08, alpha=0.18, zorder=4)
        pill(ax, sx2, 1.9, 0.05, 0.42, col, radius=0.04, alpha=1.0, zorder=5)
        T(ax, sx2+0.69, 2.11, lbl, size=11.5, bold=True, color=col,
          ha='center', zorder=5)

    # 장점 태그
    for ki, txt in enumerate(['기호만 보면 즉시 이해 — 교육 불필요',
                               '음이름(도레미) 기반 직관적 체계',
                               '앱/웹에서 공유·재생 즉시 가능']):
        py = 1.55 - ki*0.48
        ax.plot([rx+0.36],[py], 's', ms=4.5, color=hx(P['grn']), zorder=6)
        T(ax, rx+0.5, py, txt, size=10, color=P['grn'], zorder=6)

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 3  –  OMR 기법 & 모델 배경
# ─────────────────────────────────────────────────────────────────────────────
def slide_omr_background():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.4)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, 'OMR 기법과 모델 개발 배경', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, 'Optical Music Recognition — 왜 직접 학습했는가', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '04 / 16', size=9, color=P['mute'], ha='right')

    # OMR 개념 설명 (상단 배너)
    pill(ax, 0.4, 5.55, 12.5, 0.72, P['bg4'], radius=0.1, zorder=3)
    T(ax, 0.75, 5.91, 'OMR(Optical Music Recognition)이란?', size=12.5, bold=True, color=P['pri'])
    T(ax, 5.35, 5.91, '악보 이미지에서 음표·리듬·기호를 자동 인식 → 연주 가능한 데이터로 변환하는 AI 기술',
      size=11, color=P['dim'])

    # 3 컬럼 비교
    cols = [
        {
            'title': '기존 방식',
            'sub': 'Rule-based OMR',
            'color': P['mute'],
            'items': ['오선 위치 → 음표 좌표 계산', '기호별 하드코딩 규칙', '새 기호 추가 = 코드 수정'],
            'note': '유지보수 어렵고\n확장 불가',
        },
        {
            'title': '참고 모델',
            'sub': 'AndroMR / HOMR (GitHub)',
            'color': P['sky'],
            'items': ['Android 기반 오픈소스 OMR', 'C++ JNI 엔진 구조 참고', 'MusicXML 출력 방식'],
            'note': '커스텀 토큰 체계\n지원 불가',
        },
        {
            'title': '우리 모델',
            'sub': 'DeepScore (직접 학습)',
            'color': P['pri'],
            'items': ['DeepScore 토큰 1,013개', '커리큘럼 학습 R1→R4', '블록 악보 직접 출력'],
            'note': 'RTX 3080\n직접 학습',
        },
    ]
    cw = 3.9; gap = 0.28; x0 = 0.42; ch = 4.55
    for i, c in enumerate(cols):
        cx = x0 + i*(cw+gap)
        is_ours = i == 2
        fc = hx(P['bg3']) if not is_ours else hx(P['pri'])
        ec = hx(P['brh']) if not is_ours else hx(P['pdk'])
        card(ax, cx, 0.62, cw, ch, fc=fc if is_ours else (1,1,1),
             ec=ec, lw=2.0 if is_ours else 1.2, shadow=True)

        # 상단 헤더
        ax.add_patch(FancyBboxPatch((cx, 0.62+ch-0.78), cw, 0.78,
                                    boxstyle='round,pad=0.06',
                                    facecolor=hx(c['color']),
                                    edgecolor='none', zorder=5))
        tc = 'white' if i > 0 else P['dim']
        T(ax, cx+cw/2, 0.62+ch-0.37, c['title'], size=14, bold=True,
          color=tc, ha='center')
        T(ax, cx+cw/2, 0.62+ch-0.65, c['sub'], size=8.5,
          color=tc, ha='center', alpha=0.85)

        # 항목
        yi = 0.62 + ch - 1.22
        for item in c['items']:
            dot_c = c['color'] if i > 0 else P['mute']
            ax.plot([cx+0.28],[yi], 'o', ms=5, color=hx(dot_c), zorder=6)
            T(ax, cx+0.48, yi, item, size=10.5,
              color=P['txt'] if not is_ours else P['bg'], zorder=6)
            yi -= 0.62

        # 하단 노트
        pill(ax, cx+0.3, 0.72, cw-0.6, 0.52, c['color'], radius=0.07,
             alpha=0.13, zorder=5)
        T(ax, cx+cw/2, 0.98, c['note'], size=9, color=c['color'],
          bold=True, ha='center', linespacing=1.3, zorder=6)

        # 화살표
        if i < 2:
            ax.annotate('', xy=(cx+cw+gap*0.1, 0.62+ch*0.5),
                        xytext=(cx+cw+gap*0.9, 0.62+ch*0.5),
                        arrowprops=dict(arrowstyle='->', color=hx(P['pri']),
                                       lw=2.0, mutation_scale=16), zorder=7)

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 4  –  플랫폼 구성
# ─────────────────────────────────────────────────────────────────────────────
def slide_platform():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.45)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '플랫폼 구성', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '단일 OMR 파이프라인으로 웹 + 앱 서비스 제공', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '05 / 16', size=9, color=P['mute'], ha='right')

    # 플로우 박스 정의
    nodes = [
        {'label': '악보\n이미지', 'sub': 'PNG / JPG\n갤러리·카메라', 'color': P['mute'], 'x':0.38},
        {'label': 'SegNet',      'sub': 'MobileNetV3\n6-class 분류', 'color': P['sky'], 'x':3.14},
        {'label': 'Encoder',     'sub': 'ConvNeXt-T\n[B,320,512]',   'color': P['pri'], 'x':5.90},
        {'label': 'Decoder',     'sub': 'Transformer\nDeepScore 토큰','color': P['prp'], 'x':8.66},
        {'label': '표기법\n변환', 'sub': '블록 악보\n렌더링',         'color': P['grn'], 'x':11.42},
    ]
    bw=2.42; bh=2.0; by=3.1
    for nd in nodes:
        x = nd['x']
        card(ax, x, by, bw, bh, fc=(1,1,1), shadow=True, lw=1.5)
        ax.add_patch(FancyBboxPatch((x, by+bh-0.52), bw, 0.52,
                                    boxstyle='round,pad=0.05',
                                    facecolor=hx(nd['color']), edgecolor='none', zorder=5))
        T(ax, x+bw/2, by+bh-0.25, nd['label'], size=12, bold=True,
          color='white', ha='center', linespacing=1.3, zorder=6)
        T(ax, x+bw/2, by+bh*0.32, nd['sub'], size=9.5,
          color=P['dim'], ha='center', linespacing=1.5, zorder=5)
        if nd['x'] < 11.42:
            arrow(ax, x+bw+0.02, by+bh/2, nd['x']+bw+0.65, by+bh/2,
                  color=P['pri'], lw=1.8)

    # 출력 분기
    fork_x = 11.42+bw+0.25; fork_y = by+bh/2
    ax.plot([fork_x, fork_x], [fork_y-0.7, fork_y+0.7],
            color=hx(P['pri']), lw=1.5, zorder=4)
    for (lbl, sub, col, yy) in [('웹 (Web)', 'FastAPI + 가상 피아노', P['grn'], fork_y+0.7),
                                  ('앱 (Flutter)', '연습 + 공유', P['prp'], fork_y-0.7)]:
        pill(ax, fork_x+0.06, yy-0.28, 1.12, 0.56, col, radius=0.07, alpha=0.15, zorder=4)
        ax.add_patch(mpatches.Rectangle((fork_x+0.06, yy-0.28), 0.05, 0.56,
                                        facecolor=hx(col), edgecolor='none', zorder=5))
        T(ax, fork_x+0.22, yy+0.08, lbl, size=10.5, bold=True, color=col)
        T(ax, fork_x+0.22, yy-0.14, sub, size=8.5, color=P['dim'])
        ax.plot([fork_x, fork_x+0.05],[yy, yy], color=hx(col), lw=1.5, zorder=5)

    # 하단 기술 스택
    techs = ['PyTorch 2.x','TFLite INT8','ONNX Runtime','OpenCV 4.9','Flutter 3.x','FastAPI']
    x_t = 0.5
    for t in techs:
        pill(ax, x_t, 0.42, 1.85, 0.44, P['bg3'], radius=0.07, zorder=4)
        T(ax, x_t+0.93, 0.64, t, size=9.5, color=P['tdm' if t!='PyTorch 2.x' else 'pri'],
          ha='center', zorder=5)
        x_t += 2.1

    T(ax, 6.67, 2.72, '3단 파이프라인: 분류 → 특징 추출 → 시퀀스 디코딩', size=10,
      color=P['dim'], ha='center')

    return fig2stream(fig)

P['tdm'] = P['dim']
P['sdk'] = P['sky']

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6  –  추론 알고리즘 상세
# ─────────────────────────────────────────────────────────────────────────────
def slide_inference_detail():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.38)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '추론 알고리즘 상세', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55,
      'SegNet → Staff 검출 → Encoder (ConvNeXt-T) → Decoder (Transformer) → DeepScore 출력',
      size=10.5, color=P['dim'])
    T(ax, 12.9, 6.55, '06 / 16', size=9, color=P['mute'], ha='right')

    # ── 왼쪽 열: 전처리 → SegNet → Staff 검출 ───────────────────────────
    lx, cw = 0.35, 5.85
    card(ax, lx, 0.38, cw, 5.72, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, lx, 0.38, 5.72, P['sky'], w=0.07)
    T(ax, lx+0.35, 5.78, 'Phase 1  전처리 & 픽셀 분류', size=11, bold=True, color=P['sky'])
    ax.plot([lx+0.3, lx+cw-0.15],[5.62,5.62], color=hx(P['brd']), lw=0.8)

    steps_l = [
        ('전처리',    P['mute'],
         ['autocrop → resize (308px 높이)',
          'color_adjust → noise_filtering',
          '입력: PNG/JPG  출력: 정규화 텐서']),
        ('SegNet',    P['sky'],
         ['MobileNetV3-Small + U-Net 디코더',
          '6클래스 픽셀 분류 (음표/오선/쉼표 등)',
          '슬라이딩 윈도우 256×256 / stride 128',
          '23% 패치만 악보 내용 포함 (개선 예정)']),
        ('Staff 검출',P['pri'],
         ['수평 투영 합산 → Gaussian smoothing',
          'NMS 피크 검출 → 5선 그루핑',
          '제약: peaks < 5 이면 빈 배열 반환',
          '캔버스 생성: 256×1280 (단일) / 384×1280 (Grand)']),
    ]
    yi = 5.38
    for stage, col, details in steps_l:
        pill(ax, lx+0.28, yi-0.9, cw-0.5, 0.28, col, radius=0.05, alpha=1.0, zorder=5)
        T(ax, lx+0.52, yi-0.76, stage, size=10.5, bold=True, color='white', zorder=6)
        for di, d in enumerate(details):
            T(ax, lx+0.52, yi-1.18-di*0.36, d, size=9.5, color=P['dim'], zorder=5)
        yi -= 0.9 + 0.3 + len(details)*0.36 + 0.22

    # ── 오른쪽 열: Encoder → Decoder → 출력 ─────────────────────────────
    rx = 7.13
    card(ax, rx, 0.38, cw, 5.72, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, rx, 0.38, 5.72, P['prp'], w=0.07)
    T(ax, rx+0.35, 5.78, 'Phase 2  시퀀스 인식 & 출력', size=11, bold=True, color=P['prp'])
    ax.plot([rx+0.3, rx+cw-0.15],[5.62,5.62], color=hx(P['brd']), lw=0.8)

    steps_r = [
        ('Encoder',   P['pri'],
         ['ConvNeXt-Tiny 백본',
          '출력 특징: [B, 320, 512]',
          'LR = base_lr / 1.5 (워밍업 후 해제)',
          'TFLite INT8 변환 후 모바일 배포']),
        ('Decoder',   P['prp'],
         ['Autoregressive Transformer (8층)',
          '어휘 1,013 토큰 (rhythm/pitch/lift/note)',
          'KV-cache 32개 텐서, MAX_SEQ=608',
          '현재: Greedy argmax  →  개선: Beam=4']),
        ('출력',      P['grn'],
         ['DeepScore 토큰 → rhythm_rules 적용',
          'accidental_rules → MusicXML 생성',
          '커스텀 블록 악보 렌더링 (Flutter/Web)']),
    ]
    yi2 = 5.38
    for stage, col, details in steps_r:
        pill(ax, rx+0.28, yi2-0.9, cw-0.5, 0.28, col, radius=0.05, alpha=1.0, zorder=5)
        T(ax, rx+0.52, yi2-0.76, stage, size=10.5, bold=True, color='white', zorder=6)
        for di, d in enumerate(details):
            T(ax, rx+0.52, yi2-1.18-di*0.36, d, size=9.5, color=P['dim'], zorder=5)
        yi2 -= 0.9 + 0.3 + len(details)*0.36 + 0.22

    # 중앙 구분선 + 화살표
    mx = (lx+cw+rx) / 2
    ax.plot([mx, mx],[0.5, 6.1], color=hx(P['brd']), lw=1.0, zorder=2, linestyle='--')
    arrow(ax, lx+cw+0.06, 3.5, rx-0.06, 3.5, color=P['pri'], lw=2.0, ms=14)
    T(ax, mx, 3.28, '특징 맵 전달', size=9, color=P['mute'], ha='center')

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 7  –  진행 단계 (4그룹)
# ─────────────────────────────────────────────────────────────────────────────
def slide_phases():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.45)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '현재 진행 단계', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '비슷한 성격의 단계를 4개 그룹으로 통합', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '07 / 16', size=9, color=P['mute'], ha='right')

    groups = [
        {
            'num': '01', 'title': '준비',
            'phases': ['Phase 0  표기법 규칙 확정', 'Phase 1  데이터 생성 파이프라인'],
            'color': P['grn'], 'status': '완료', 'icon': '✓',
        },
        {
            'num': '02', 'title': '학습 & 변환',
            'phases': ['Phase 2  AI 모델 학습 (R1~R3)', 'Phase 3  TFLite INT8 변환'],
            'color': P['pri'], 'status': '진행 중', 'icon': '◉',
        },
        {
            'num': '03', 'title': '엔진 & UI',
            'phases': ['Phase 4  C++ 추론 엔진', 'Phase 5  Flutter 앱', 'Phase 6  커스텀 렌더러'],
            'color': P['sky'], 'status': '대부분 완료', 'icon': '◐',
        },
        {
            'num': '04', 'title': '배포 & 서비스',
            'phases': ['Phase 7  서버 배포', 'Phase 8  웹 OMR 연결', 'Phase 9  커뮤니티'],
            'color': P['prp'], 'status': '예정', 'icon': '○',
        },
    ]

    n = 4; cw = 2.98; gap = 0.18; cy_timeline = 5.15
    for i, g in enumerate(groups):
        cx = 0.38 + i*(cw+gap)
        cc = hx(g['color'])
        is_done  = g['status'] == '완료'
        is_todo  = g['status'] == '예정'

        # 타임라인 연결선
        if i > 0:
            lc = hx(P['grn']) if groups[i-1]['status']=='완료' else hx(P['brd'])
            ax.plot([cx-gap, cx], [cy_timeline, cy_timeline],
                    color=lc, lw=2.5, zorder=2)

        # 원형 아이콘
        circle = plt.Circle((cx+cw/2, cy_timeline), 0.38,
                             color=cc, zorder=5)
        ax.add_patch(circle)
        T(ax, cx+cw/2, cy_timeline, g['icon'], size=15, bold=True,
          color='white', ha='center', zorder=6)

        # 번호 + 제목
        T(ax, cx+cw/2, cy_timeline+0.78, g['num'], size=9.5, bold=True,
          color=g['color'], ha='center', alpha=0.7)
        T(ax, cx+cw/2, cy_timeline+0.48, g['title'], size=13.5, bold=True,
          color=g['color'], ha='center')

        # 카드 (하단)
        card_top = 1.22; card_h = 3.55
        card(ax, cx, card_top, cw, card_h,
             fc=hx(P['bg3']) if not is_todo else (1,1,1),
             ec=hx(g['color']), lw=1.8 if not is_todo else 1.0,
             shadow=not is_todo)

        yp = card_top + card_h - 0.52
        for ph in g['phases']:
            ax.plot([cx+0.32],[yp], 's', ms=5,
                    color=cc if not is_todo else hx(P['mute']),
                    alpha=0.7, zorder=6)
            T(ax, cx+0.50, yp, ph, size=9.5,
              color=P['txt'] if not is_todo else P['mute'], zorder=6)
            yp -= 0.56

        # 상태 배지
        pill(ax, cx+0.35, card_top+0.14, cw-0.7, 0.38,
             g['color'], radius=0.06, alpha=0.12 if is_todo else 0.18, zorder=5)
        T(ax, cx+cw/2, card_top+0.33, g['status'], size=10,
          bold=True, color=g['color'], ha='center', zorder=6)

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 6  –  웹페이지
# ─────────────────────────────────────────────────────────────────────────────
def slide_webpage():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.4)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '웹페이지 (online_webpage)', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '가상 피아노 + 커스텀 악보 뷰어 + OMR 업로드', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '08 / 16', size=9, color=P['mute'], ha='right')

    # 브라우저 목업 (왼쪽)
    card(ax, 0.38, 0.45, 7.05, 5.72, fc=(1,1,1), shadow=True, lw=1.0)

    # 브라우저 탭바
    ax.add_patch(mpatches.Rectangle((0.38, 5.78), 7.05, 0.39,
                                    facecolor=hx(P['bg2']), edgecolor='none', zorder=5))
    for xc, col in [(0.56,P['red']),(0.73,P['amb']),(0.90,P['grn'])]:
        ax.plot([xc],[5.97], 'o', ms=6, color=hx(col), zorder=6)
    T(ax, 3.9, 5.97, 'musicscore-demo.vercel.app', size=9, color=P['dim'], ha='center')

    # 내비게이션
    ax.add_patch(mpatches.Rectangle((0.38, 5.38), 7.05, 0.4,
                                    facecolor=hx(P['pri']), edgecolor='none', zorder=5))
    for ni, nav in enumerate(['튜토리얼','악보 변환','연주하기']):
        T(ax, 1.26+ni*2.1, 5.58, nav, size=10, bold=True, color='white', ha='center', zorder=6)

    # 콘텐츠 영역: 샘플 악보
    T(ax, 3.9, 5.2, '샘플 악보 선택', size=9, color=P['mute'], ha='center')
    n_col = [P['pri'],P['sky'],P['grn'],'#F5A623',P['amb'],P['prp']]
    n_lbl = ['C','D','E','F','G','A']
    xl = 0.55
    for nc, nl in zip(n_col, n_lbl):
        pill(ax, xl, 4.52, 1.0, 0.56, nc, radius=0.07, zorder=6)
        T(ax, xl+0.5, 4.80, nl, size=16, bold=True, color='white', ha='center', zorder=7)
        xl += 1.08

    # 박자 블록 예시
    T(ax, 3.9, 4.28, '표기법 렌더링 결과', size=9, color=P['mute'], ha='center')
    dur_row = [('C',1,P['pri']),('E',0.5,P['grn']),('G',0.5,P['amb']),
               ('A',1,P['prp']),('B',2,P['sky'])]
    total_u = 5; xs=0.52; unit_=6.0/total_u
    for nn_, d_, nc_ in dur_row:
        bw_ = d_*unit_-0.05
        pill(ax, xs, 3.6, bw_, 0.56, nc_, radius=0.06, zorder=6)
        T(ax, xs+bw_/2, 3.88, nn_, size=13 if d_>=1 else 10,
          bold=True, color='white', ha='center', zorder=7)
        xs += d_*unit_

    # 가상 피아노 (간략)
    T(ax, 3.9, 3.35, '88건반 가상 피아노', size=9, color=P['mute'], ha='center')
    kw_ = 6.6/14; n14 = 14
    for ki in range(n14):
        pill(ax, 0.52+ki*kw_, 1.1, kw_-0.03, 2.1,
             P['bg'] if ki < 7 else P['bg2'], radius=0.04, zorder=5)
        ax.add_patch(mpatches.Rectangle((0.52+ki*kw_, 1.1), kw_-0.03, 0.06,
                                        facecolor=hx(P['brd']), edgecolor='none', zorder=6))
    bk_= {1,2,4,5,6,8,9,11,12}
    for ki in bk_:
        if ki < n14:
            pill(ax, 0.52+ki*kw_+kw_*0.62, 1.78, kw_*0.44, 1.42,
                 P['txt'], radius=0.03, zorder=7)

    T(ax, 3.9, 0.76, '키보드 단축키 + 마우스 클릭 지원', size=8.5,
      color=P['mute'], ha='center')

    # 오른쪽: 완료 / 예정 기능
    card(ax, 7.72, 3.35, 5.2, 2.82, fc=(1,1,1), shadow=True)
    accent_bar(ax, 7.72, 3.35, 2.82, P['grn'], w=0.07)
    T(ax, 8.08, 5.75, '구현 완료', size=10, bold=True, color=P['grn'])
    ax.plot([8.08,12.68],[5.57,5.57], color=hx(P['brd']), lw=0.7)
    done_list = ['샘플 악보 → 블록 표기법 즉시 렌더링',
                 '88건반 가상 피아노 (키보드 + 마우스)',
                 '튜토리얼 (기호 읽는 법 안내)',
                 '옥타브 연동 + 그라데이션 테마 UI']
    yd = 5.35
    for item in done_list:
        ax.plot([8.08],[yd],'o', ms=4.5, color=hx(P['grn']), zorder=5)
        T(ax, 8.26, yd, item, size=10.5, color=P['txt'])
        yd -= 0.48

    card(ax, 7.72, 0.45, 5.2, 2.7, fc=(1,1,1), shadow=True)
    accent_bar(ax, 7.72, 0.45, 2.7, P['amb'], w=0.07)
    T(ax, 8.08, 2.82, '개발 예정', size=10, bold=True, color=P['amb'])
    ax.plot([8.08,12.68],[2.65,2.65], color=hx(P['brd']), lw=0.7)
    todo_list = ['악보 이미지 업로드 → OMR API 변환',
                 'FastAPI 서버 연동',
                 'MIDI 키보드 지원 (Web MIDI API)']
    yt = 2.42
    for item in todo_list:
        ax.plot([8.08],[yt],'o', ms=4.5, color=hx(P['amb']), zorder=5)
        T(ax, 8.26, yt, item, size=10.5, color=P['txt'])
        yt -= 0.58

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 9  –  학습 정확도 현황 그래프
# ─────────────────────────────────────────────────────────────────────────────
def slide_accuracy_graph():
    # ── 데이터 ─────────────────────────────────────────────────────────────
    #  rows = Round 1 / Round 2 / Round 3
    #  cols = Phase1(SegNet) / Phase2(Enc+Dec) / Phase3(E2E)
    data = [
        [98.70, 65.73, 70.94],  # Round 1
        [98.60, 44.97, 44.61],  # Round 2
        [None,  48.76, 48.76],  # Round 3
    ]
    col_labels = ['Phase 1\nSegNet', 'Phase 2\nEncoder+Decoder', 'Phase 3\nEnd-to-End']
    row_labels = ['Round 1', 'Round 2', 'Round 3']
    col_colors = [P['grn'], P['sky'], P['pri']]

    def acc_color(v):
        if v is None: return P['brd']
        if v >= 90:   return P['grn']
        if v >= 60:   return P['sky']
        if v >= 45:   return P['amb']
        return P['red']

    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.3)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '학습 정확도 현황', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, 'val_acc = max(0, 1 − TER)  |  TER: Token Error Rate (Levenshtein / GT 길이)',
      size=10.5, color=P['dim'])
    T(ax, 12.9, 6.55, '09 / 16', size=9, color=P['mute'], ha='right')

    # ── 왼쪽: 색상 테이블 ──────────────────────────────────────────────────
    tx0 = 0.35;  tw = 7.6
    th_row = 0.80; th_hdr = 0.55
    col_w  = tw / 3
    table_top = 5.95

    # 컬럼 헤더
    for ci, (lbl, col) in enumerate(zip(col_labels, col_colors)):
        cx = tx0 + ci*col_w
        ax.add_patch(FancyBboxPatch((cx+0.06, table_top-th_hdr+0.06), col_w-0.12, th_hdr-0.06,
                                    boxstyle='round,pad=0.04',
                                    facecolor=hx(col), edgecolor='none', zorder=4))
        T(ax, cx+col_w/2, table_top-th_hdr/2+0.03, lbl, size=10, bold=True,
          color='white', ha='center', linespacing=1.35, zorder=5)

    # 행 헤더 + 데이터 셀
    for ri, (row_lbl, row_data) in enumerate(zip(row_labels, data)):
        ry = table_top - th_hdr - (ri+1)*th_row
        # 행 라벨
        pill(ax, tx0-0.02, ry+0.06, 0.95, th_row-0.12, P['bg4'],
             radius=0.06, zorder=3)
        T(ax, tx0+0.44, ry+th_row/2, row_lbl, size=11, bold=True,
          color=P['txt'], ha='center', zorder=4)

        # 데이터 셀
        for ci, val in enumerate(row_data):
            cx  = tx0 + ci*col_w + 0.04
            col = acc_color(val)
            # 셀 배경
            ax.add_patch(FancyBboxPatch((tx0 + ci*col_w + 0.06, ry+0.07),
                                        col_w-0.12, th_row-0.14,
                                        boxstyle='round,pad=0.05',
                                        facecolor=hx(col), edgecolor='none',
                                        alpha=0.15, zorder=3))
            # 좌측 색상 바
            ax.add_patch(mpatches.Rectangle(
                (tx0 + ci*col_w+0.06, ry+0.07), 0.05, th_row-0.14,
                facecolor=hx(col), edgecolor='none', zorder=4))
            # 수치
            txt = f'{val:.2f}%' if val is not None else '—'
            fs  = 16 if val is not None else 18
            T(ax, tx0 + ci*col_w + col_w/2 + 0.03, ry+th_row/2,
              txt, size=fs, bold=True, color=col if val is not None else P['mute'],
              ha='center', zorder=5)

    # 행 구분선
    for ri in range(len(data)+1):
        ry = table_top - th_hdr - ri*th_row
        ax.plot([tx0, tx0+tw],[ry, ry], color=hx(P['brd']), lw=0.8, zorder=2)

    # 범례 (색상 의미)
    legend_items = [('≥ 90%', P['grn']), ('≥ 60%', P['sky']),
                    ('≥ 45%', P['amb']), ('< 45%', P['red'])]
    lx = tx0 + 0.05
    T(ax, lx, 0.78, '정확도 등급:', size=9, color=P['dim'], bold=True)
    lx += 1.05
    for ltxt, lcol in legend_items:
        pill(ax, lx, 0.6, 0.85, 0.34, lcol, radius=0.05, alpha=0.85, zorder=4)
        T(ax, lx+0.43, 0.77, ltxt, size=9, bold=True, color='white', ha='center', zorder=5)
        lx += 1.0

    # ── 오른쪽: 꺾은선 그래프 ──────────────────────────────────────────────
    ax2 = fig.add_axes([0.615, 0.11, 0.365, 0.76])
    ax2.set_facecolor((1,1,1,0))

    rnd_x = [1, 2, 3]
    e2e = [70.94, 44.61, 48.76]
    p2  = [65.73, 44.97, 48.76]

    ax2.plot(rnd_x, [98.7, 98.6, None], 'o-', color=hx(P['grn']),
             lw=2.2, ms=8, label='Phase 1  SegNet', clip_on=False)
    ax2.plot(rnd_x, p2, 's--', color=hx(P['sky']),
             lw=2.0, ms=8, label='Phase 2  Enc+Dec', clip_on=False)
    ax2.plot(rnd_x, e2e, 'D-', color=hx(P['pri']),
             lw=2.2, ms=8, label='Phase 3  E2E', clip_on=False)

    # 수치 레이블
    for xi, (s, p, e) in enumerate(zip([98.7,98.6],[65.73,44.97,48.76][0:2],e2e)):
        ax2.text(xi+1+0.06, s+1.5, f'{s:.1f}', fontsize=8.5,
                 fontproperties=fm.FontProperties(fname=_KRB), color=hx(P['grn']))
    for xi, val in enumerate(p2):
        ax2.text(xi+1+0.06, val+1.5, f'{val:.1f}', fontsize=8.5,
                 fontproperties=fm.FontProperties(fname=_KRB), color=hx(P['sky']))
    for xi, val in enumerate(e2e):
        ax2.text(xi+1-0.18, val-4.5, f'{val:.2f}', fontsize=8.5,
                 fontproperties=fm.FontProperties(fname=_KRB), color=hx(P['pri']))

    ax2.axhline(72.01, color=hx(P['mute']), lw=1.0, ls='--', alpha=0.7)
    ax2.text(3.06, 73, '이전 R1 최고\n72.01%', fontsize=7.5,
             fontproperties=fm.FontProperties(fname=_KR), color=hx(P['mute']), va='bottom')

    ax2.set_xticks(rnd_x)
    ax2.set_xticklabels(['R1', 'R2', 'R3'],
                        fontproperties=fm.FontProperties(fname=_KR), fontsize=11)
    ax2.set_ylim(0, 108); ax2.set_xlim(0.6, 3.8)
    ax2.set_ylabel('val_acc (%)', fontproperties=fm.FontProperties(fname=_KR), fontsize=9)
    ax2.set_title('라운드별 정확도 추이', fontproperties=fm.FontProperties(fname=_KRB),
                  fontsize=11, color=hx(P['txt']), pad=6)
    ax2.legend(prop=fm.FontProperties(fname=_KR), fontsize=8,
               loc='upper right', framealpha=0.85, edgecolor=hx(P['brd']))
    ax2.spines[['top','right']].set_visible(False)
    ax2.set_axisbelow(True)
    ax2.grid(axis='y', alpha=0.25, lw=0.7)
    ax2.yaxis.set_tick_params(labelsize=9)

    # 구분선
    ax.plot([8.1, 8.1],[0.35, 6.32], color=hx(P['brd']), lw=1.0, ls='--', zorder=2)

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 10  –  핵심 문제 상세
# ─────────────────────────────────────────────────────────────────────────────
def slide_issues_detail():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.35)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '핵심 문제 상세', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '현재 모델의 4가지 병목 — 원인 분석 및 수치 근거', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '10 / 16', size=9, color=P['mute'], ha='right')

    issues = [
        {
            'num': '①', 'title': 'SegNet 패치 효율 23%',
            'col': P['amb'],
            'cause': '원인: A4 이미지(1240×1754) 중 악보 영역은 상단 13.5%뿐\n'
                     '      랜덤 패치 샘플링 → 77%가 순백색 배경',
            'fix':   '해결: y = 89~330px 범위 집중 샘플링 → 효율 3~4배 향상',
            'code':  'y0 = random.randint(0, H-P)  # 완전 랜덤 (현재)',
        },
        {
            'num': '②', 'title': '그리디 디코딩 연쇄 오류',
            'col': P['red'],
            'cause': '원인: argmax 1개 토큰 오류 → 이후 시퀀스 전체 붕괴\n'
                     '      R2 희귀 토큰(꾸밈음·셋잇단음) 특히 취약',
            'fix':   '해결: Beam Search (beam=4) 교체 → 후보군 유지로 오류 전파 억제',
            'code':  'next_id = int(np.argmax(logits[0,0]))  # line 478',
        },
        {
            'num': '③', 'title': '소규모 데이터 학습 한계',
            'col': P['pri'],
            'cause': '원인: 200장 × batch=8 = 25 steps/epoch × 3 epoch = 75 steps\n'
                     '      어휘 1,013 토큰 커버에 최소 1,500~2,000장 필요',
            'fix':   '해결: 처음부터 충분한 데이터 확보 후 학습 시작 (점진 추가 X)',
            'code':  'MAX_VAL = 200  # val samples / epoch (train.py)',
        },
        {
            'num': '④', 'title': 'Grand Staff 그루핑 취약',
            'col': P['prp'],
            'cause': '원인: 두 오선 사이 간격이 MAX_UNIT 이하이면\n'
                     '      잘못된 5선 그루핑 (gap이 내부 간격과 혼동)',
            'fix':   '해결: 이중모드 간격 분석 → treble/bass 자동 구분 알고리즘',
            'code':  'if all(MIN_UNIT<=g<=MAX_UNIT for g in gaps)  # line 308',
        },
    ]

    col_xs = [0.35, 6.84]
    for idx, iss in enumerate(issues):
        cx = col_xs[idx % 2]
        cy = 3.05 if idx < 2 else 0.38
        cw2, ch2 = 6.14, 2.52
        col = iss['col']

        card(ax, cx, cy, cw2, ch2, fc=(1,1,1), shadow=True, lw=1.4)
        pill(ax, cx, cy+ch2-0.38, cw2, 0.38, col, radius=0.05, alpha=1.0, zorder=5)
        T(ax, cx+0.22, cy+ch2-0.19, f"{iss['num']}  {iss['title']}",
          size=11.5, bold=True, color='white', zorder=6)

        T(ax, cx+0.22, cy+ch2-0.72, iss['cause'], size=9.5,
          color=P['txt'], linespacing=1.5, zorder=5)
        pill(ax, cx+0.18, cy+ch2-1.55, cw2-0.36, 0.32, col, alpha=0.08, zorder=4)
        T(ax, cx+0.25, cy+ch2-1.39, iss['fix'], size=9.5, bold=True,
          color=col, zorder=5)
        pill(ax, cx+0.18, cy+0.1, cw2-0.36, 0.38, P['bg4'], alpha=1.0, zorder=4)
        T(ax, cx+0.28, cy+0.29, iss['code'], size=8.5,
          color=P['dim'], zorder=5)

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 11  –  개선 방향 & 참고
# ─────────────────────────────────────────────────────────────────────────────
# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 11  –  Round 1 재학습 결과
# ─────────────────────────────────────────────────────────────────────────────
def slide_r1_retraining():
    # 비교 데이터
    phases = ['Phase 1\nSegNet', 'Phase 2\nEnc+Dec', 'Phase 3\nEnd-to-End']
    new_vals = [98.70, 65.73, 70.94]   # ROOT 새 학습 (2026.06.22)
    old_vals = [98.70, 65.73, 72.01]   # 이전 최고 R1
    col_ph   = [P['grn'], P['sky'], P['pri']]

    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.32)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, 'Round 1 재학습 결과', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '2026.06.22  ROOT 가중치 재측정  —  이전 최고치와 단계별 비교',
      size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '11 / 16', size=9, color=P['mute'], ha='right')

    # ── 왼쪽: 단계별 비교 카드 3개 ──────────────────────────────────────
    for pi, (ph, nv, ov, col) in enumerate(zip(phases, new_vals, old_vals, col_ph)):
        cx = 0.35 + pi*3.08; cw2 = 2.82; cy = 1.62; ch = 4.45
        card(ax, cx, cy, cw2, ch, fc=(1,1,1), shadow=True, lw=1.4)
        # 컬러 헤더
        ax.add_patch(FancyBboxPatch((cx, cy+ch-0.50), cw2, 0.50,
                                    boxstyle='round,pad=0.05',
                                    facecolor=hx(col), edgecolor='none', zorder=5))
        T(ax, cx+cw2/2, cy+ch-0.25, ph, size=10.5, bold=True,
          color='white', ha='center', linespacing=1.3, zorder=6)

        diff = nv - ov
        diff_col = P['grn'] if diff >= 0 else P['red']
        diff_txt = f'+{diff:.2f}%p' if diff >= 0 else f'{diff:.2f}%p'

        # 새 학습 수치 (크게)
        T(ax, cx+cw2/2, cy+ch-1.18, '신규 학습', size=9, color=P['mute'],
          ha='center', zorder=5)
        T(ax, cx+cw2/2, cy+ch-1.72, f'{nv:.2f}%', size=24, bold=True,
          color=col, ha='center', zorder=6)

        # 구분선
        ax.plot([cx+0.25, cx+cw2-0.25],[cy+ch-2.1, cy+ch-2.1],
                color=hx(P['brd']), lw=0.8, zorder=4)

        # 이전 최고 수치
        T(ax, cx+cw2/2, cy+ch-2.42, '이전 최고', size=9, color=P['mute'],
          ha='center', zorder=5)
        T(ax, cx+cw2/2, cy+ch-2.92, f'{ov:.2f}%', size=18, bold=False,
          color=P['dim'], ha='center', zorder=5)

        # 차이 배지
        pill(ax, cx+cw2/2-0.6, cy+0.78, 1.2, 0.44, diff_col,
             radius=0.07, alpha=0.15, zorder=4)
        pill(ax, cx+cw2/2-0.6, cy+0.78, 0.05, 0.44, diff_col,
             radius=0.04, alpha=1.0, zorder=5)
        T(ax, cx+cw2/2+0.08, cy+1.0, diff_txt, size=13, bold=True,
          color=diff_col, ha='center', zorder=5)
        T(ax, cx+cw2/2, cy+0.52, '차이', size=8.5, color=P['mute'], ha='center', zorder=4)

    # ── 오른쪽: 해석 & 결론 ──────────────────────────────────────────────
    rx = 9.62; rw = 3.36
    card(ax, rx, 1.62, rw, 4.45, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, rx, 1.62, 4.45, P['pri'], w=0.07)
    T(ax, rx+0.35, 5.78, '결과 해석', size=11, bold=True, color=P['pri'])
    ax.plot([rx+0.3, rx+rw-0.15],[5.62, 5.62], color=hx(P['brd']), lw=0.8)

    insights = [
        (P['grn'], 'SegNet 동등 수준 유지',
         '98.70% — 픽셀 분류 안정적'),
        (P['sky'], 'Phase 2 실질적 동등',
         '65.73% — 동일 수렴점 도달'),
        (P['amb'], 'Phase 3 소폭 하락',
         '70.94% vs 72.01% (−1.07%p)\n랜덤 시드 차이 허용 범위'),
        (P['grn'], 'R2 시작 가능 판단',
         'ROOT 가중치 = 공식 R1 기준\n→ ml/models/round1/ 복사 완료'),
    ]
    yi = 5.38
    for col, title, desc in insights:
        ax.plot([rx+0.35],[yi], 'o', ms=7, color=hx(col),
                markerfacecolor=hx(col), zorder=6)
        T(ax, rx+0.58, yi, title, size=10.5, bold=True, color=P['txt'], zorder=5)
        T(ax, rx+0.58, yi-0.38, desc, size=9, color=P['dim'],
          linespacing=1.4, zorder=5)
        yi -= 1.02

    # ── 하단: 학습 파라미터 요약 배너 ────────────────────────────────────
    pill(ax, 0.35, 0.28, 12.63, 1.18, P['bg4'], radius=0.08, zorder=3)
    T(ax, 0.62, 1.22, '학습 환경', size=9.5, bold=True, color=P['pri'])
    params = [
        ('GPU', 'RTX 3080'),
        ('Framework', 'PyTorch 2.x + OneCycleLR + AdamW'),
        ('Phase 1 epoch', '50  /  batch 16  /  lr 1e-4'),
        ('Phase 2 epoch', '100  /  batch 8  /  lr 1e-4'),
        ('Phase 3 epoch', '50  /  batch 4  /  lr 3e-5'),
        ('데이터', 'Round 1 합성 악보  ~200장'),
    ]
    px = 0.62; py = 0.72
    for ki, (k, v) in enumerate(params):
        T(ax, px, py, f'{k}:', size=8.5, color=P['dim'], bold=True)
        T(ax, px+1.2, py, v, size=8.5, color=P['txt'])
        px += 2.14
        if ki == 2:
            px = 0.62; py -= 0.4

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 12  –  개선 방향 & 참고
# ─────────────────────────────────────────────────────────────────────────────
def slide_improvements_detail():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.35)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '개선 방향 & 참고 모델', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '단계별 구현 계획 + 오픈소스 참고 알고리즘', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '12 / 16', size=9, color=P['mute'], ha='right')

    # ── 상단 3개 개선 카드 ───────────────────────────────────────────────
    fixes = [
        {
            'tag': '즉시 적용', 'col': P['grn'],
            'title': 'SegNet 패치 집중 샘플링',
            'steps': ['dataset.py line 424-425 수정',
                      'y0 = random.randint(89, min(330, H-P))',
                      '예상 효율: 77% → 23% 빈 패치로 감소'],
            'effect': '학습 효율 3~4배 향상 기대',
        },
        {
            'tag': 'R2 전 적용', 'col': P['sky'],
            'title': 'Beam Search (beam=4) 교체',
            'steps': ['omr_inference.py line 478 교체',
                      'train.py greedy_decode() 동시 수정',
                      '빔 크기 4 → 메모리 4× 증가 허용 범위'],
            'effect': '희귀 토큰 연쇄 오류 대폭 감소',
        },
        {
            'tag': 'R4 전 적용', 'col': P['pri'],
            'title': 'IMSLP 실사진 + 증강 파이프라인',
            'steps': ['IMSLP 공개 악보 직접 촬영 50~100장',
                      'OpenCV: 원근왜곡·노이즈·JPEG 아티팩트',
                      '목표: 1,500~2,000장 확보'],
            'effect': '실환경 인식률 70%→85%+ 목표',
        },
    ]
    for fi, fx in enumerate(fixes):
        cx = 0.35 + fi*4.34
        cw3, ch3 = 4.0, 3.35
        card(ax, cx, 2.72, cw3, ch3, fc=(1,1,1), shadow=True, lw=1.4)
        pill(ax, cx, 2.72+ch3-0.36, cw3, 0.36, fx['col'], radius=0.05, zorder=5)
        pill(ax, cx+0.2, 2.72+ch3-0.65, 0.82, 0.22, fx['col'], alpha=0.18, zorder=5)
        T(ax, cx+0.61, 2.72+ch3-0.54, fx['tag'], size=8.5, bold=True,
          color=fx['col'], ha='center', zorder=6)
        T(ax, cx+0.22, 2.72+ch3-0.19, fx['title'], size=11, bold=True,
          color='white', zorder=6)
        yi = 2.72+ch3-0.88
        for step in fx['steps']:
            ax.plot([cx+0.3],[yi], 'o', ms=3.5, color=hx(fx['col']), zorder=5)
            T(ax, cx+0.45, yi, step, size=9.5, color=P['dim'], zorder=5)
            yi -= 0.45
        pill(ax, cx+0.18, 2.82, cw3-0.36, 0.36, fx['col'], alpha=0.12, zorder=4)
        T(ax, cx+cw3/2, 3.0, fx['effect'], size=9.5, bold=True,
          color=fx['col'], ha='center', zorder=5)

    # ── 하단: 참고 모델 ──────────────────────────────────────────────────
    ax.plot([0.35, 12.98],[2.6, 2.6], color=hx(P['brd']), lw=1.0)
    T(ax, 0.45, 2.42, '참고 오픈소스 모델', size=11.5, bold=True, color=P['txt'])

    refs = [
        {
            'name': 'SMT++', 'col': P['prp'],
            'meta': 'ICDAR 2024  |  MIT License  |  arxiv 2407.xxxxx',
            'desc': ['멀티시스템 전 페이지 OMR', '그랜드 스태프 처리 구조 참고',
                     '인코더: DeiT-B / 디코더: Transformer'],
        },
        {
            'name': 'LEGATO', 'col': P['sky'],
            'meta': '2025.06  |  arxiv 2506.19065  |  공개 예정',
            'desc': ['214K 이미지, ABC 표기법', '전 페이지 다중 시스템 인식',
                     '토큰 체계 설계 참고 가능'],
        },
        {
            'name': 'Grand Staff\n개선 알고리즘', 'col': P['grn'],
            'meta': '내부 개발 계획 (R3 전)',
            'desc': ['이중 모드 간격 분석 (treble/bass 구분)',
                     'Bimodal gap detection',
                     '오선 그루핑 오류율 0% 목표'],
        },
    ]
    for ri, ref in enumerate(refs):
        rx2 = 0.35 + ri*4.34
        card(ax, rx2, 0.38, 4.0, 2.0, fc=(1,1,1), shadow=True, lw=1.2)
        accent_bar(ax, rx2, 0.38, 2.0, ref['col'], w=0.06)
        T(ax, rx2+0.28, 2.15, ref['name'], size=11, bold=True, color=ref['col'])
        T(ax, rx2+0.28, 1.84, ref['meta'], size=8.5, color=P['mute'])
        ax.plot([rx2+0.22, rx2+3.85],[1.7,1.7], color=hx(P['brd']), lw=0.6)
        for di, d in enumerate(ref['desc']):
            ax.plot([rx2+0.32],[1.52-di*0.36],'o',ms=3.5, color=hx(ref['col']),zorder=5)
            T(ax, rx2+0.46, 1.52-di*0.36, d, size=9.5, color=P['dim'])

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 12  –  AI 모델 파이프라인 상세
# ─────────────────────────────────────────────────────────────────────────────
def slide_pipeline():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.4)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, 'AI 모델 파이프라인', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, 'SegNet → Encoder → Decoder  3단계 구조', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '13 / 16', size=9, color=P['mute'], ha='right')

    stages = [
        {'title': '입력', 'detail': '악보 이미지\nPNG / JPG',
         'tech': 'OpenCV 4.9\n전처리',
         'color': P['mute'], 'badge': None},
        {'title': 'SegNet', 'detail': 'MobileNetV3\n+ U-Net 디코더\n6-class 픽셀 분류\nFocal Loss',
         'tech': 'val_acc  98.9%',
         'color': P['sky'], 'badge': ('Phase 1', P['grn'])},
        {'title': 'Encoder', 'detail': 'ConvNeXt-Tiny\n백본\n→ [B, 320, 512]',
         'tech': 'AMP + AdamW',
         'color': P['pri'], 'badge': ('Phase 2', P['pri'])},
        {'title': 'Decoder', 'detail': 'Autoregressive\nTransformer\nTeacher Forcing\nLabel Smooth 0.1',
         'tech': 'vocab 1,013',
         'color': P['prp'], 'badge': ('Phase 2/3', P['prp'])},
        {'title': '출력', 'detail': 'DeepScore 토큰\n→ 블록 악보\n렌더링',
         'tech': '커스텀 표기법',
         'color': P['grn'], 'badge': None},
    ]
    bw=2.32; bh=4.05; by=1.52; gap_=0.28
    total_w = len(stages)*bw + (len(stages)-1)*gap_
    x0 = (13.33-total_w)/2

    for i, s in enumerate(stages):
        cx = x0 + i*(bw+gap_)
        card(ax, cx, by, bw, bh, fc=(1,1,1), shadow=True, lw=1.5)
        # 헤더
        ax.add_patch(FancyBboxPatch((cx, by+bh-0.58), bw, 0.58,
                                    boxstyle='round,pad=0.05',
                                    facecolor=hx(s['color']), edgecolor='none', zorder=5))
        T(ax, cx+bw/2, by+bh-0.27, s['title'], size=13.5, bold=True,
          color='white' if s['color']!=P['mute'] else P['dim'],
          ha='center', zorder=6)

        # 상세
        T(ax, cx+bw/2, by+bh*0.55, s['detail'], size=9.5, color=P['dim'],
          ha='center', linespacing=1.5, zorder=5)

        # 하단 tech 뱃지
        pill(ax, cx+0.22, by+0.2, bw-0.44, 0.42, s['color'],
             radius=0.06, alpha=0.13, zorder=5)
        T(ax, cx+bw/2, by+0.41, s['tech'], size=9, bold=True,
          color=s['color'], ha='center', zorder=6)

        # Phase 뱃지 (위)
        if s['badge']:
            bl, bc = s['badge']
            pill(ax, cx+bw/2-0.65, by-0.52, 1.3, 0.35, bc, radius=0.06, alpha=0.15, zorder=5)
            T(ax, cx+bw/2, by-0.35, bl, size=8.5, bold=True, color=bc, ha='center', zorder=6)

        # 화살표
        if i < len(stages)-1:
            arrow(ax, cx+bw+0.02, by+bh/2, cx+bw+gap_-0.02, by+bh/2,
                  color=P['pri'], lw=2)

    # 하단 설명
    T(ax, 6.67, 1.12, 'PyTorch 2.x  +  AMP  |  OneCycleLR  |  RTX 3080  10 GB  |  배치 8~16',
      size=9.5, color=P['mute'], ha='center')

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 8  –  커리큘럼 학습
# ─────────────────────────────────────────────────────────────────────────────
def slide_curriculum():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.4)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '커리큘럼 학습 전략', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '난이도 단계별 누적 학습 — 각 Round는 이전 가중치를 이어받아 파인튜닝', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '14 / 16', size=9, color=P['mute'], ha='right')

    rounds = [
        {'rnd':'R1','title':'기초','color':P['pri'],'data':'합성 2,000장','status':'완료',
         'items':['음표 / 쉼표','박자표 / 조표','도돌이표']},
        {'rnd':'R2','title':'심화','color':P['sky'],'data':'합성 2,000장','status':'재학습 중',
         'items':['셈여림 / 헤어핀','아티큘레이션','꾸밈음 / 이음줄','셋잇단음표 / 8va']},
        {'rnd':'R3','title':'Grand Staff','color':P['amb'],'data':'합성 2,000장','status':'예정',
         'items':['2단 보표 (피아노)','Treble + Bass 통합']},
        {'rnd':'R4','title':'실사진','color':P['grn'],'data':'사진 1,000장+','status':'예정',
         'items':['실제 악보 사진','도메인 적응 파인튜닝']},
    ]
    n=4; cw=2.92; gap_=0.28
    total_w = n*cw+(n-1)*gap_
    x0=(13.33-total_w)/2

    for i,r in enumerate(rounds):
        cx = x0+i*(cw+gap_)
        c=r['color']; bh=4.55
        card(ax, cx, 0.9, cw, bh, fc=(1,1,1), lw=2.0, shadow=True)
        # 헤더
        ax.add_patch(FancyBboxPatch((cx, 0.9+bh-0.72), cw, 0.72,
                                    boxstyle='round,pad=0.06',
                                    facecolor=hx(c), edgecolor='none', zorder=5))
        T(ax, cx+cw/2, 0.9+bh-0.5, r['rnd'], size=9.5, color='white',
          ha='center', alpha=0.85, zorder=6)
        T(ax, cx+cw/2, 0.9+bh-0.27, r['title'], size=14, bold=True,
          color='white', ha='center', zorder=6)

        # 항목
        yi = 0.9+bh-1.2
        for item in r['items']:
            ax.plot([cx+0.3],[yi],'>', ms=5, color=hx(c), alpha=0.7, zorder=6)
            T(ax, cx+0.5, yi, item, size=10, color=P['txt'], zorder=6)
            yi -= 0.55

        # 데이터 / 상태
        ax.plot([cx+0.25, cx+cw-0.25],[1.72,1.72], color=hx(P['brd']), lw=0.7)
        T(ax, cx+cw/2, 1.55, r['data'], size=9.5, color=P['dim'], ha='center')
        pill(ax, cx+0.28, 0.98, cw-0.56, 0.42, c, radius=0.06, alpha=0.15, zorder=5)
        T(ax, cx+cw/2, 1.19, r['status'], size=9.5, bold=True, color=c,
          ha='center', zorder=6)

        # 화살표
        if i < n-1:
            arrow(ax, cx+cw+0.02, 0.9+bh/2, cx+cw+gap_-0.02, 0.9+bh/2,
                  color=P['pri'], lw=1.8, ms=13)

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 9  –  학습 결과
# ─────────────────────────────────────────────────────────────────────────────
def slide_results():
    fig = plt.figure(figsize=(13.33, 7.5), facecolor='white')

    # 헤더 (matplotlib로 직접)
    ax_hdr = fig.add_axes([0,0.852,1,0.148])
    ax_hdr.axis('off'); ax_hdr.set_xlim(0,1); ax_hdr.set_ylim(0,1)
    ax_hdr.set_facecolor('white')
    ax_hdr.plot([0,1],[0,0], color=hx(P['pri']), lw=3, transform=ax_hdr.transAxes,
                clip_on=False)
    fp_b = fm.FontProperties(fname=_KRB)
    fp_r = fm.FontProperties(fname=_KR)
    ax_hdr.text(0.034, 0.72, '학습 결과 현황', fontsize=22, fontproperties=fp_b,
                color=hx(P['txt']), va='center')
    ax_hdr.text(0.034, 0.26, '코드 버그 4개 수정 완료 → 재학습 예정', fontsize=11,
                fontproperties=fp_r, color=hx(P['dim']), va='center')
    ax_hdr.text(0.968, 0.26, '15 / 16', fontsize=9, fontproperties=fp_r,
                color=hx(P['mute']), va='center', ha='right')

    # 배경 그라데이션
    ax_bg = fig.add_axes([0,0,1,0.852])
    ax_bg.axis('off')
    gradient_bg(ax_bg, alpha=0.4)

    gs = GridSpec(1, 3, figure=fig,
                  left=0.04, right=0.97, top=0.82, bottom=0.08,
                  wspace=0.08)

    # ── SegNet 차트 ────────────────────────────────────────────────────────────
    ax1 = fig.add_subplot(gs[0])
    ax1.set_facecolor(hx(P['bg2']))
    bars1 = ax1.bar(['Round 1', 'Round 2'], [98.9, 98.6],
                    color=[hx(P['grn']), hx(P['sky'])],
                    width=0.45, zorder=3)
    ax1.set_ylim(95, 101)
    ax1.set_title('SegNet  픽셀 분류 정확도',
                  fontproperties=fm.FontProperties(fname=_KRB),
                  fontsize=12, color=hx(P['txt']), pad=10)
    ax1.set_ylabel('val_acc (%)', fontproperties=fm.FontProperties(fname=_KR),
                   fontsize=10, color=hx(P['dim']))
    for sp in ax1.spines.values():
        sp.set_color(hx(P['brh'])); sp.set_linewidth(0.6)
    ax1.tick_params(colors=hx(P['dim']), labelsize=10)
    ax1.yaxis.label.set_fontproperties(fm.FontProperties(fname=_KR))
    ax1.grid(axis='y', color=hx(P['brd']), lw=0.6, zorder=0)
    for bar, v in zip(bars1, [98.9, 98.6]):
        ax1.text(bar.get_x()+bar.get_width()/2, v+0.04, f'{v}%',
                 ha='center', va='bottom', fontsize=12, fontweight='bold',
                 color=hx(P['txt']),
                 fontproperties=fm.FontProperties(fname=_KRB))
    ax1.text(0.5, 0.05, '안정 수렴 — 병목 아님',
             transform=ax1.transAxes, ha='center', fontsize=9.5,
             fontproperties=fp_r,
             color=hx(P['grn']),
             bbox=dict(fc=hx(P['grn']), alpha=0.1,
                       boxstyle='round,pad=0.3', ec=hx(P['grn'])))

    # ── Seq2Seq 차트 ───────────────────────────────────────────────────────────
    ax2 = fig.add_subplot(gs[1])
    ax2.set_facecolor(hx(P['bg2']))
    cats = ['R1\nval', 'R1 공식\n(200)', 'R1 공식\n(500)', 'R2\nval']
    vals = [72, 84.4, 91.7, 64]
    clrs = [hx(P['pri']), hx(P['grn']), hx(P['sky']), hx(P['amb'])]
    bars2 = ax2.bar(cats, vals, color=clrs, width=0.45, zorder=3)
    ax2.set_ylim(0, 112)
    ax2.set_title('Seq2Seq  토큰 인식 정확도',
                  fontproperties=fm.FontProperties(fname=_KRB),
                  fontsize=12, color=hx(P['txt']), pad=10)
    ax2.set_ylabel('Accuracy (%)', fontproperties=fp_r,
                   fontsize=10, color=hx(P['dim']))
    for sp in ax2.spines.values():
        sp.set_color(hx(P['brh'])); sp.set_linewidth(0.6)
    ax2.tick_params(colors=hx(P['dim']), labelsize=10)
    ax2.grid(axis='y', color=hx(P['brd']), lw=0.6, zorder=0)
    for bar, v in zip(bars2, vals):
        ax2.text(bar.get_x()+bar.get_width()/2, v+0.8, f'{v}%',
                 ha='center', va='bottom', fontsize=11, fontweight='bold',
                 color=hx(P['txt']),
                 fontproperties=fp_b)
    ax2.annotate('R2 하락\n재학습 필요', xy=(3, 64), xytext=(2.25, 82),
                 fontsize=8.5, color=hx(P['amb']), ha='center',
                 fontproperties=fp_r,
                 arrowprops=dict(arrowstyle='->', color=hx(P['amb']),
                                 lw=1.2, mutation_scale=10))

    # ── 버그 수정 목록 ─────────────────────────────────────────────────────────
    ax3 = fig.add_subplot(gs[2])
    ax3.axis('off'); ax3.set_facecolor('white')
    ax3.text(0.5, 0.96, '수정된 버그 4가지', ha='center', va='top', fontsize=12,
             fontproperties=fp_b, color=hx(P['txt']))
    ax3.axhline(0.88, color=hx(P['brd']), lw=0.8)

    bugs = [
        ('①', 'Grand Staff 레이블 불일치',
         'is_system 플래그 + 5-tuple 수정', P['pri']),
        ('②', '검증 샘플 부족 (50 → 200)',
         'MAX_VAL = 200 으로 상향', P['sky']),
        ('③', 'Encoder 학습률 너무 낮음',
         'LR 나눗수  3.0 → 1.5 로 조정', P['grn']),
        ('④', '평가 전처리 불일치',
         'detect_staffs + extract_canvas 적용', P['amb']),
    ]
    yb = 0.82
    for num, title, detail, col in bugs:
        ax3.add_patch(FancyBboxPatch((0.04, yb-0.17), 0.92, 0.19,
                                     boxstyle='round,pad=0.02',
                                     facecolor=hx(col), alpha=0.08,
                                     edgecolor=hx(col), linewidth=0.8,
                                     transform=ax3.transAxes))
        ax3.add_patch(mpatches.Rectangle((0.04, yb-0.17), 0.025, 0.19,
                                          facecolor=hx(col), edgecolor='none',
                                          transform=ax3.transAxes))
        ax3.text(0.09, yb-0.04, f'{num} {title}', fontsize=10, fontproperties=fp_b,
                 color=hx(P['txt']), va='center', transform=ax3.transAxes)
        ax3.text(0.09, yb-0.13, detail, fontsize=9, fontproperties=fp_r,
                 color=hx(P['dim']), va='center', transform=ax3.transAxes)
        yb -= 0.22

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 10  –  앞으로 할 것 + 일정
# ─────────────────────────────────────────────────────────────────────────────
def slide_next():
    import datetime
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.4)

    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '앞으로 할 것 & 일정', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '2026년 8월 17일 데모 목표', size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '16 / 16', size=9, color=P['mute'], ha='right')

    # 3 단계 카드
    steps = [
        ('단기   6월 ~ 7월 초', P['pri'], [
            'Round 1 재학습  (목표 val 80%+)',
            'Round 2, 3 누적 학습',
            'detect_staffs 성공률 확인',
        ]),
        ('중기   7월', P['sky'], [
            'TFLite INT8 변환 완료',
            'FastAPI OMR 서버 배포',
            'Flutter 렌더러 구현',
        ]),
        ('데모   8월', P['grn'], [
            '웹 OMR API 연결 완성',
            '실기기 QA 완료',
            '영상 / 포스터 제출',
        ]),
    ]
    sw=3.85; gap_s=0.27
    total_sw = 3*sw+2*gap_s; x0s=(13.33-total_sw)/2
    for i,(title,col,items) in enumerate(steps):
        cx=x0s+i*(sw+gap_s)
        card(ax, cx, 3.92, sw, 2.22, fc=(1,1,1), shadow=True)
        ax.add_patch(mpatches.Rectangle((cx, 3.92+2.22-0.5), sw, 0.5,
                                        facecolor=hx(col), edgecolor='none', zorder=5))
        T(ax, cx+sw/2, 3.92+2.22-0.24, title, size=11.5, bold=True,
          color='white', ha='center', zorder=6)
        yi = 3.92+2.22-0.82
        for item in items:
            ax.plot([cx+0.28],[yi],'>', ms=4.5, color=hx(col), alpha=0.7, zorder=6)
            T(ax, cx+0.46, yi, item, size=10, color=P['txt'], zorder=6)
            yi -= 0.5

    # Gantt
    ax_g = fig.add_axes([0.03, 0.07, 0.94, 0.44])
    ax_g.set_facecolor(hx(P['bg2']))
    fp_b2 = fm.FontProperties(fname=_KRB)
    fp_r2 = fm.FontProperties(fname=_KR)

    base = datetime.date(2026,6,13); end_d = datetime.date(2026,8,17)
    total_d = (end_d-base).days

    tasks = [
        ('Round 1 재학습',      '6/13','6/23', P['pri']),
        ('Round 2 재학습',      '6/23','7/8',  P['sky']),
        ('Round 3 학습',        '6/24','7/8',  P['amb']),
        ('TFLite + 서버',       '7/9', '7/24', P['prp']),
        ('Flutter 렌더러',      '7/9', '7/24', P['grn']),
        ('앱 / 웹 QA',          '7/25','8/4',  P['sky']),
        ('영상 / 포스터',        '8/5', '8/15', P['red']),
    ]

    def pd(s):
        m,d=s.split('/'); return datetime.date(2026,int(m),int(d))

    n_t = len(tasks)
    for i,(name,ts,te,col) in enumerate(tasks):
        xs=(pd(ts)-base).days/total_d
        xe=(pd(te)-base).days/total_d
        y=n_t-1-i
        ax_g.barh(y, xe-xs, left=xs, height=0.52,
                  color=hx(col), alpha=0.82, zorder=3)
        if xe-xs > 0.05:
            ax_g.text((xs+xe)/2, y, name, ha='center', va='center',
                      fontsize=8.5, color='white', fontweight='bold',
                      fontproperties=fp_b2, zorder=4)

    for d_str, lbl in [('6/13','6/13'),('7/1','7/1'),('8/1','8/1')]:
        xp=(pd(d_str)-base).days/total_d
        ax_g.axvline(xp, color=hx(P['brh']), lw=0.8, alpha=0.7)
        ax_g.text(xp+0.005, n_t+0.1, lbl, fontsize=8,
                  color=hx(P['dim']), fontproperties=fp_r2)

    ax_g.axvline(1.0, color=hx(P['red']), lw=2, linestyle='--', zorder=5)
    ax_g.text(0.995, n_t+0.1, '8/17 데모', fontsize=8.5, ha='right',
              color=hx(P['red']), fontweight='bold', fontproperties=fp_b2)

    ax_g.set_xlim(0, 1.05); ax_g.set_ylim(-0.5, n_t+0.55)
    ax_g.set_xticks([]); ax_g.set_yticks([])
    for sp in ax_g.spines.values():
        sp.set_color(hx(P['brh'])); sp.set_linewidth(0.4)
    ax_g.grid(axis='x', color=hx(P['brd']), lw=0.4, alpha=0.6)

    return fig2stream(fig)

# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE 2  –  현황 분석 요약  (대화 내용 기반)
# ─────────────────────────────────────────────────────────────────────────────
def slide_summary():
    fig = plt.figure(figsize=(13.33, 7.5))
    ax  = fig.add_axes([0,0,1,1])
    ax.set_xlim(0,13.33); ax.set_ylim(0,7.5); ax.axis('off')
    ax.set_facecolor('white')
    gradient_bg(ax, W=13.33, H=7.5, alpha=0.38)

    # 헤더
    ax.add_patch(mpatches.Rectangle((0,6.4),13.33,1.1,
                                    facecolor='white',edgecolor='none',zorder=2))
    ax.plot([0,13.33],[6.4,6.4], color=hx(P['pri']), lw=2.5, zorder=3)
    T(ax, 0.45, 6.95, '현황 분석 요약', size=22, bold=True, color=P['txt'])
    T(ax, 0.45, 6.55, '맞춤형 악보 인식 & 변환 앱  —  2026 캡스톤 디자인  |  분석 기준: 2026.06.22',
      size=11, color=P['dim'])
    T(ax, 12.9, 6.55, '02 / 16', size=9, color=P['mute'], ha='right')

    # ── 열 1: 학습 정확도 현황 ───────────────────────────────────────────────
    col_x = [0.35, 4.65, 8.95]
    cw    = 4.0

    card(ax, col_x[0], 0.42, cw, 5.72, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, col_x[0], 0.42, 5.72, P['pri'], w=0.07)
    T(ax, col_x[0]+0.35, 5.82, '학습 정확도 현황', size=11, bold=True, color=P['pri'])
    ax.plot([col_x[0]+0.3, col_x[0]+cw-0.15],[5.65,5.65],
            color=hx(P['brd']), lw=0.8)

    # SegNet 섹션
    T(ax, col_x[0]+0.3, 5.42, 'SegNet  (픽셀 분류)', size=9.5,
      bold=True, color=P['dim'])
    rows_seg = [('Round 1', '98.9%', P['grn']),
                ('Round 2', '98.6%', P['grn'])]
    ys = 5.12
    for lbl, val, col in rows_seg:
        pill(ax, col_x[0]+0.3, ys-0.18, cw-0.6, 0.34,
             col, radius=0.05, alpha=0.1, zorder=4)
        T(ax, col_x[0]+0.48, ys, lbl, size=10, color=P['dim'])
        T(ax, col_x[0]+cw-0.25, ys, val, size=12, bold=True,
          color=col, ha='right')
        ys -= 0.44

    ax.plot([col_x[0]+0.3, col_x[0]+cw-0.15],[ys+0.24, ys+0.24],
            color=hx(P['brd']), lw=0.6)
    ys -= 0.06

    # Seq2Seq 섹션
    T(ax, col_x[0]+0.3, ys, 'Seq2Seq  (토큰 인식)', size=9.5,
      bold=True, color=P['dim'])
    ys -= 0.34
    rows_seq = [('Round 1  (코드 수정 후)', '70.94%', P['pri']),
                ('Round 1  (이전 최고치)',  '72.01%', P['mute']),
                ('Round 2',               '65.74%', P['amb']),
                ('Round 3',               '48.76%', P['red'])]
    for lbl, val, col in rows_seq:
        pill(ax, col_x[0]+0.3, ys-0.17, cw-0.6, 0.32,
             col, radius=0.05, alpha=0.08, zorder=4)
        T(ax, col_x[0]+0.48, ys, lbl, size=9.5, color=P['dim'])
        T(ax, col_x[0]+cw-0.25, ys, val, size=11.5, bold=True,
          color=col, ha='right')
        ys -= 0.42

    ax.plot([col_x[0]+0.3, col_x[0]+cw-0.15],[ys+0.22, ys+0.22],
            color=hx(P['brd']), lw=0.6)
    ys -= 0.06

    # 앱/웹
    T(ax, col_x[0]+0.3, ys, '앱 & 웹', size=9.5, bold=True, color=P['dim'])
    ys -= 0.34
    for lbl, val, col in [('Flutter 앱 UI', '완료', P['grn']),
                           ('웹 가상 피아노', '완료', P['grn']),
                           ('OMR API 연결', '예정', P['amb'])]:
        T(ax, col_x[0]+0.48, ys, lbl, size=9.5, color=P['dim'])
        pill(ax, col_x[0]+cw-0.95, ys-0.14, 0.75, 0.3,
             col, radius=0.05, alpha=0.15, zorder=4)
        T(ax, col_x[0]+cw-0.57, ys, val, size=9.5, bold=True,
          color=col, ha='center')
        ys -= 0.38

    # ── 열 2: 오늘 발견한 핵심 문제 ─────────────────────────────────────────
    card(ax, col_x[1], 0.42, cw, 5.72, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, col_x[1], 0.42, 5.72, P['red'], w=0.07)
    T(ax, col_x[1]+0.35, 5.82, '발견된 핵심 문제', size=11, bold=True, color=P['red'])
    ax.plot([col_x[1]+0.3, col_x[1]+cw-0.15],[5.65,5.65],
            color=hx(P['brd']), lw=0.8)

    issues = [
        ('①', '훈련 이미지 빈공간 98.6%',
         'SegNet 패치의 77%가 순백\n→ 악보 기호 학습 효율 저하',
         P['amb']),
        ('②', '그리디 디코딩의 한계',
         '틀린 토큰 1개 → 이후 시퀀스 연쇄 오류\n→ R2 희귀 토큰(꾸밈음, 셋잇단) 특히 취약',
         P['red']),
        ('③', '소규모 데이터 학습 효과 없음',
         '200장 → acc 32%  /  점진적 추가도 역효과\n→ 처음부터 1,500~2,000장 필요',
         P['pri']),
        ('④', 'Grand Staff 그루핑 로직 취약',
         '두 오선 사이 간격이 MAX_UNIT 이하면\n잘못된 5선 그루핑 발생 가능',
         P['prp']),
    ]
    yi = 5.38
    for num, title, desc, col in issues:
        pill(ax, col_x[1]+0.28, yi-0.88, cw-0.52, 1.06,
             col, radius=0.08, alpha=0.06, zorder=3)
        ax.add_patch(mpatches.Rectangle(
            (col_x[1]+0.28, yi-0.88), 0.05, 1.06,
            facecolor=hx(col), edgecolor='none', zorder=4))
        T(ax, col_x[1]+0.52, yi-0.18, f'{num}  {title}',
          size=11, bold=True, color=P['txt'])
        T(ax, col_x[1]+0.52, yi-0.6, desc, size=9.5,
          color=P['dim'], linespacing=1.5)
        yi -= 1.28

    # ── 열 3: 개선 방향 & 참고 정보 ─────────────────────────────────────────
    card(ax, col_x[2], 0.42, cw, 5.72, fc=(1,1,1), shadow=True, lw=1.4)
    accent_bar(ax, col_x[2], 0.42, 5.72, P['grn'], w=0.07)
    T(ax, col_x[2]+0.35, 5.82, '개선 방향 & 참고', size=11, bold=True, color=P['grn'])
    ax.plot([col_x[2]+0.3, col_x[2]+cw-0.15],[5.65,5.65],
            color=hx(P['brd']), lw=0.8)

    fixes = [
        ('즉시', [
            'SegNet 패치: 악보 영역 집중 샘플링',
            '(빈공간 77% → 학습 효율 3배 향상 기대)',
        ], P['pri']),
        ('R2 전', [
            '그리디 → 빔 서치 (beam = 4) 교체',
            'omr_inference.py + train.py 동시 수정',
        ], P['sky']),
        ('R4 전', [
            'IMSLP 공개 악보 직접 촬영 50~100장',
            'OpenCV 증강: 원근왜곡·노이즈·JPEG 아티팩트',
        ], P['grn']),
        ('참고 모델', [
            'SMT++  (MIT 라이선스, 멀티시스템)',
            'LEGATO  (2025.06, 전 페이지 OMR)',
        ], P['prp']),
    ]
    yi2 = 5.38
    for tag, items, col in fixes:
        pill(ax, col_x[2]+0.28, yi2-0.88, cw-0.52, 1.06,
             col, radius=0.08, alpha=0.06, zorder=3)
        pill(ax, col_x[2]+0.28, yi2-0.3, 0.72, 0.28,
             col, radius=0.05, alpha=0.2, zorder=4)
        T(ax, col_x[2]+0.64, yi2-0.16, tag, size=9, bold=True,
          color=col, ha='center')
        for ki, item in enumerate(items):
            T(ax, col_x[2]+0.52, yi2-0.56-ki*0.32, item,
              size=9.5, color=P['dim' if ki else 'txt'])
        yi2 -= 1.28

    # ── 하단 라이선스 배너 ────────────────────────────────────────────────────
    pill(ax, 0.35, 0.12, 12.63, 0.28, P['bg4'], radius=0.06, zorder=3)
    T(ax, 0.7, 0.26,
      '라이선스 검토 완료:  PyTorch BSD-3  /  ONNX Runtime MIT  /  OpenCV Apache 2.0  '
      '/  TFLite Apache 2.0  /  tinyxml2 zlib  →  전 스택 상업 이용 가능',
      size=9, color=P['pri'], bold=True)

    return fig2stream(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  ASSEMBLE PPTX
# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_fullslide(stream):
    sl = prs.slides.add_slide(BLANK)
    sl.shapes.add_picture(stream, Inches(0), Inches(0),
                          prs.slide_width, prs.slide_height)

print('슬라이드 생성 중...')
builders = [
    ('타이틀',          slide_title),
    ('현황 분석 요약',   slide_summary),
    ('프로젝트 소개',    slide_overview),
    ('OMR 배경',        slide_omr_background),
    ('플랫폼',           slide_platform),
    ('추론 알고리즘',    slide_inference_detail),
    ('진행 단계',        slide_phases),
    ('웹페이지',         slide_webpage),
    ('학습 정확도 그래프', slide_accuracy_graph),
    ('핵심 문제 상세',   slide_issues_detail),
    ('R1 재학습 결과',   slide_r1_retraining),
    ('개선 방향 & 참고', slide_improvements_detail),
    ('AI 파이프라인',    slide_pipeline),
    ('커리큘럼',         slide_curriculum),
    ('학습 결과',        slide_results),
    ('앞으로 + 일정',    slide_next),
]
for name, fn in builders:
    print(f'  {name}...')
    add_fullslide(fn())

out = Path(__file__).resolve().parent.parent.parent / 'project_status_v4.pptx'
prs.save(str(out))
print(f'\n완료: {out}')
